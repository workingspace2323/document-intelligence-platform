import fitz
import pytesseract
from PIL import Image
import io
import docx
import pandas as pd
from pptx import Presentation

pytesseract.pytesseract.tesseract_cmd = r"C:\Users\Admin\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"


# ---------------- PDF (OCR Supported) ---------------- #

def read_pdf(path):

    doc = fitz.open(path)

    full_text = ""
    pages_data = []

    for i in range(len(doc)):

        page = doc[i]

        text = page.get_text().strip()

        if not text:

            pix = page.get_pixmap()

            img = Image.open(
                io.BytesIO(
                    pix.tobytes("png")
                )
            )

            text = pytesseract.image_to_string(
                img
            )

        full_text += text + "\n"

        pages_data.append(
            {
                "page": i + 1,
                "text": text
            }
        )

    return {
        "text": full_text,
        "pages": len(doc),
        "characters": len(full_text),
        "pages_data": pages_data
    }


# ---------------- DOCX ---------------- #

def read_docx(path):

    document = docx.Document(path)

    text = "\n".join(
        [p.text for p in document.paragraphs]
    )

    return {
        "text": text,
        "pages": 1,
        "characters": len(text),
        "pages_data": [
            {
                "page": 1,
                "text": text
            }
        ]
    }


# ---------------- TXT ---------------- #

def read_txt(path):

    with open(
        path,
        "r",
        encoding="utf-8",
        errors="ignore"
    ) as f:

        text = f.read()

    return {
        "text": text,
        "pages": 1,
        "characters": len(text),
        "pages_data": [
            {
                "page": 1,
                "text": text
            }
        ]
    }


# ---------------- EXCEL ---------------- #

def read_excel(path):

    df = pd.read_excel(path)

    text = df.astype(str).to_string()

    return {
        "text": text,
        "pages": 1,
        "characters": len(text),
        "pages_data": [
            {
                "page": 1,
                "text": text
            }
        ]
    }


# ---------------- PPTX ---------------- #

def read_pptx(path):

    presentation = Presentation(path)

    text = ""

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return {
        "text": text,
        "pages": len(presentation.slides),
        "characters": len(text),
        "pages_data": [
            {
                "page": 1,
                "text": text
            }
        ]
    }


# ---------------- MAIN ROUTER ---------------- #

def extract_file(path):

    extension = path.lower().split(".")[-1]

    if extension == "pdf":
        return read_pdf(path)

    elif extension == "docx":
        return read_docx(path)

    elif extension == "txt":
        return read_txt(path)

    elif extension in ["xlsx", "xls"]:
        return read_excel(path)

    elif extension == "pptx":
        return read_pptx(path)

    else:
        raise ValueError(
            f"Unsupported file type: {extension}"
        )